from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# Theme
BG_LIGHT = RGBColor(248, 250, 252)
TEXT_DARK = RGBColor(15, 23, 42)
TEXT_MUTED = RGBColor(71, 85, 105)
ACCENT = RGBColor(245, 158, 11)
ACCENT_DARK = RGBColor(217, 119, 6)
GREEN = RGBColor(16, 185, 129)
BLUE = RGBColor(59, 130, 246)
RED = RGBColor(239, 68, 68)
BORDER = RGBColor(203, 213, 225)


def set_run_style(run, size=14, bold=False, color=TEXT_DARK, font_name="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def add_full_bg(slide, color=BG_LIGHT):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_header(slide, title, subtitle):
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.05))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(255, 255, 255)
    banner.line.fill.background()

    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.0), Inches(13.333), Inches(0.05))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.2), Inches(8.5), Inches(0.45))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run_style(run, size=25, bold=True, color=TEXT_DARK, font_name="Calibri")

    subtitle_box = slide.shapes.add_textbox(Inches(0.67), Inches(0.67), Inches(8.8), Inches(0.28))
    p2 = subtitle_box.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = subtitle
    set_run_style(run2, size=12, bold=False, color=TEXT_MUTED, font_name="Calibri")

    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.55), Inches(0.23), Inches(2.2), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(255, 247, 237)
    badge.line.color.rgb = ACCENT
    badge.line.width = Pt(1.2)
    tf = badge.text_frame
    tf.clear()
    p3 = tf.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "BudgetWise"
    set_run_style(run3, size=12, bold=True, color=ACCENT_DARK, font_name="Calibri")


def add_card(slide, x, y, w, h, title, body, fill=RGBColor(255, 255, 255), title_color=TEXT_DARK):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = BORDER
    card.line.width = Pt(1.1)
    card.shadow.inherit = False

    tf = card.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.word_wrap = True

    p_title = tf.paragraphs[0]
    r_title = p_title.add_run()
    r_title.text = title
    set_run_style(r_title, size=14, bold=True, color=title_color, font_name="Calibri")

    p_body = tf.add_paragraph()
    r_body = p_body.add_run()
    r_body.text = body
    set_run_style(r_body, size=11, bold=False, color=TEXT_MUTED, font_name="Calibri")
    p_body.space_before = Pt(8)
    p_body.line_spacing = 1.2
    return card


def add_arrow(slide, shape_type, x, y, w, h, color=ACCENT):
    arrow = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow


def add_slide_number(slide, n):
    tb = slide.shapes.add_textbox(Inches(12.7), Inches(7.15), Inches(0.5), Inches(0.2))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(n)
    set_run_style(run, size=10, bold=False, color=TEXT_MUTED)


def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, RGBColor(255, 255, 255))

    top_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(2.15))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = RGBColor(255, 247, 237)
    top_band.line.fill.background()

    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.07), Inches(13.333), Inches(0.08))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()

    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), Inches(0.35), Inches(1.6), Inches(1.6))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = RGBColor(255, 237, 213)
    circle1.line.fill.background()

    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.45), Inches(0.95), Inches(1.2), Inches(1.2))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(254, 215, 170)
    circle2.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.6), Inches(9.5), Inches(1.2))
    t = title_box.text_frame
    p = t.paragraphs[0]
    r = p.add_run()
    r.text = "BudgetWise Project Analysis"
    set_run_style(r, size=42, bold=True, color=TEXT_DARK)

    p2 = t.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Architecture, Modules, Data Flow, Use Cases, and ER Model"
    set_run_style(r2, size=18, bold=False, color=TEXT_MUTED)
    p2.space_before = Pt(10)

    info = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5.2), Inches(7.5), Inches(1.2))
    info.fill.solid()
    info.fill.fore_color.rgb = BG_LIGHT
    info.line.color.rgb = BORDER
    info.line.width = Pt(1)

    tf = info.text_frame
    tf.clear()
    p3 = tf.paragraphs[0]
    r3 = p3.add_run()
    r3.text = "Prepared from codebase review of React + Express + MongoDB implementation."
    set_run_style(r3, size=14, bold=False, color=TEXT_DARK)

    p4 = tf.add_paragraph()
    r4 = p4.add_run()
    r4.text = "Includes implementation-aligned diagrams and six functional modules."
    set_run_style(r4, size=13, bold=False, color=TEXT_MUTED)
    p4.space_before = Pt(8)

    date_tb = slide.shapes.add_textbox(Inches(0.95), Inches(6.8), Inches(4.5), Inches(0.4))
    dp = date_tb.text_frame.paragraphs[0]
    dr = dp.add_run()
    dr.text = "Generated: April 7, 2026"
    set_run_style(dr, size=12, bold=False, color=TEXT_MUTED)


def create_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    add_header(
        slide,
        "Architecture Diagram",
        "Frontend, API, and data layers as implemented in the BudgetWise repository",
    )

    add_card(
        slide,
        0.55,
        2.4,
        1.85,
        1.55,
        "User",
        "Browser client\nTriggers dashboard,\nbudget and expense actions",
        fill=RGBColor(255, 255, 255),
    )

    add_card(
        slide,
        2.8,
        2.05,
        2.55,
        2.3,
        "React Frontend (Vite)",
        "- React Router pages\n- Shared components\n- Recharts visualizations\n- Axios API layer",
        fill=RGBColor(255, 255, 255),
    )

    add_card(
        slide,
        5.75,
        1.5,
        3.15,
        3.45,
        "Node.js + Express API",
        "- /api/budgets\n- /api/expenses\n- /api/stats\n- Input checks + errors\n- CORS and JSON middleware",
        fill=RGBColor(255, 255, 255),
    )

    add_card(
        slide,
        9.25,
        2.18,
        1.55,
        1.25,
        "Mongoose",
        "Schema +\nODM layer",
        fill=RGBColor(255, 255, 255),
    )

    add_card(
        slide,
        11.1,
        1.88,
        1.65,
        1.85,
        "MongoDB",
        "Collections:\nBudget\nExpense",
        fill=RGBColor(255, 255, 255),
    )

    add_arrow(slide, MSO_SHAPE.RIGHT_ARROW, 2.42, 2.95, 0.34, 0.33, ACCENT)
    add_arrow(slide, MSO_SHAPE.RIGHT_ARROW, 5.38, 2.95, 0.32, 0.33, ACCENT)
    add_arrow(slide, MSO_SHAPE.RIGHT_ARROW, 8.95, 2.62, 0.26, 0.3, ACCENT)
    add_arrow(slide, MSO_SHAPE.RIGHT_ARROW, 10.83, 2.62, 0.24, 0.3, ACCENT)

    return_arrow = add_arrow(slide, MSO_SHAPE.LEFT_ARROW, 5.43, 4.08, 2.58, 0.29, BLUE)
    label_box = slide.shapes.add_textbox(Inches(6.05), Inches(4.08), Inches(1.45), Inches(0.26))
    label_p = label_box.text_frame.paragraphs[0]
    label_p.alignment = PP_ALIGN.CENTER
    label_run = label_p.add_run()
    label_run.text = "JSON response"
    set_run_style(label_run, size=9, bold=True, color=RGBColor(255, 255, 255))
    return_arrow.fill.fore_color.rgb = BLUE

    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(5.45), Inches(12.0), Inches(1.25))
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(255, 255, 255)
    note.line.color.rgb = BORDER
    note.line.width = Pt(1)
    ntf = note.text_frame
    ntf.clear()
    p = ntf.paragraphs[0]
    r = p.add_run()
    r.text = "Key implementation traits:"
    set_run_style(r, size=12, bold=True, color=TEXT_DARK)
    p2 = ntf.add_paragraph()
    p2.text = "- Frontend uses /api proxy through Vite.\n- Backend seeds default budgets and expenses on first startup.\n- Stats routes compute monthly aggregates from expense documents."
    for run in p2.runs:
        set_run_style(run, size=11, bold=False, color=TEXT_MUTED)
    p2.space_before = Pt(6)
    p2.line_spacing = 1.2


def create_modules_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    add_header(
        slide,
        "Project Broken Into 6 Modules",
        "Functional decomposition of the existing codebase",
    )

    modules = [
        (
            "1) Presentation and Navigation",
            "React pages, routing, navigation shell, reusable UI cards,\nspinners, toasts, and table/chart rendering.",
            RGBColor(255, 255, 255),
        ),
        (
            "2) Budget Management",
            "Budget loading, category utilization computation, single and\nbatch limit updates, and budget health indicators.",
            RGBColor(255, 255, 255),
        ),
        (
            "3) Expense Management",
            "Expense CRUD with modal forms, category/date filtering,\ntransaction table, and inline edit/delete flows.",
            RGBColor(255, 255, 255),
        ),
        (
            "4) Analytics and Reporting",
            "Summary KPIs, 6-month trends, category breakdown,\nrecommendations, and alert generation for over-budget states.",
            RGBColor(255, 255, 255),
        ),
        (
            "5) API and Business Logic",
            "Express route handlers, request validation, pagination,\naggregation pipelines, error handling, and health endpoint.",
            RGBColor(255, 255, 255),
        ),
        (
            "6) Data Persistence Layer",
            "Mongoose models, MongoDB collections, indexes,\nseed data strategy, and schema-level constraints.",
            RGBColor(255, 255, 255),
        ),
    ]

    x_positions = [0.6, 4.55, 8.5]
    y_positions = [1.5, 4.15]
    idx = 0
    for y in y_positions:
        for x in x_positions:
            title, body, color = modules[idx]
            add_card(slide, x, y, 4.2, 2.2, title, body, fill=color)
            idx += 1


def create_data_flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    add_header(
        slide,
        "Data Flow Diagram",
        "Request-response path for budget and expense operations",
    )

    steps = [
        (0.55, 1.7, "1. User action\nAdd/Edit expense\nor update budget"),
        (3.1, 1.7, "2. React form state\nand client checks"),
        (5.65, 1.7, "3. Axios request\nto /api endpoint"),
        (8.2, 1.7, "4. Express route\nvalidation + logic"),
        (8.2, 4.05, "5. Mongoose and\nMongoDB read/write"),
        (5.65, 4.05, "6. Aggregations\n(summary, trends,\ncategories)"),
        (3.1, 4.05, "7. JSON response\n(success + data)"),
        (0.55, 4.05, "8. UI refresh\ncards/charts/tables\nalerts and badges"),
    ]

    for x, y, label in steps:
        add_card(slide, x, y, 2.15, 1.55, "", label, fill=RGBColor(255, 255, 255))

    add_arrow(slide, MSO_SHAPE.RIGHT_ARROW, 2.75, 2.3, 0.25, 0.25, ACCENT)
    add_arrow(slide, MSO_SHAPE.RIGHT_ARROW, 5.28, 2.3, 0.25, 0.25, ACCENT)
    add_arrow(slide, MSO_SHAPE.RIGHT_ARROW, 7.83, 2.3, 0.25, 0.25, ACCENT)
    add_arrow(slide, MSO_SHAPE.DOWN_ARROW, 9.08, 3.33, 0.22, 0.52, ACCENT)
    add_arrow(slide, MSO_SHAPE.LEFT_ARROW, 7.85, 4.67, 0.25, 0.25, ACCENT)
    add_arrow(slide, MSO_SHAPE.LEFT_ARROW, 5.3, 4.67, 0.25, 0.25, ACCENT)
    add_arrow(slide, MSO_SHAPE.LEFT_ARROW, 2.75, 4.67, 0.25, 0.25, ACCENT)
    add_arrow(slide, MSO_SHAPE.UP_ARROW, 1.65, 3.35, 0.22, 0.52, BLUE)

    legend = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.75), Inches(5.55), Inches(2.25), Inches(1.05))
    legend.fill.solid()
    legend.fill.fore_color.rgb = RGBColor(255, 255, 255)
    legend.line.color.rgb = BORDER
    legend.line.width = Pt(1)
    tf = legend.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Legend"
    set_run_style(r, size=11, bold=True)
    p1 = tf.add_paragraph()
    p1.text = "Orange: request path\nBlue: UI refresh loop"
    for run in p1.runs:
        set_run_style(run, size=10, color=TEXT_MUTED)
    p1.space_before = Pt(6)


def add_use_case_oval(slide, x, y, w, h, text):
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    oval.fill.solid()
    oval.fill.fore_color.rgb = RGBColor(255, 255, 255)
    oval.line.color.rgb = BORDER
    oval.line.width = Pt(1)
    tf = oval.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    set_run_style(r, size=11, bold=True, color=TEXT_DARK)
    return oval


def create_use_case_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    add_header(
        slide,
        "Use Case Diagram",
        "Primary interactions supported by the current BudgetWise application",
    )

    boundary = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.85), Inches(1.45), Inches(9.95), Inches(5.15))
    boundary.fill.solid()
    boundary.fill.fore_color.rgb = RGBColor(255, 255, 255)
    boundary.line.color.rgb = BORDER
    boundary.line.width = Pt(1.4)
    btf = boundary.text_frame
    btf.clear()
    bp = btf.paragraphs[0]
    br = bp.add_run()
    br.text = "BudgetWise System Boundary"
    set_run_style(br, size=10, bold=True, color=TEXT_MUTED)

    actor = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(2.6), Inches(1.95), Inches(1.45))
    actor.fill.solid()
    actor.fill.fore_color.rgb = RGBColor(255, 255, 255)
    actor.line.color.rgb = ACCENT
    actor.line.width = Pt(1.3)
    atf = actor.text_frame
    atf.clear()
    ap = atf.paragraphs[0]
    ap.alignment = PP_ALIGN.CENTER
    ar = ap.add_run()
    ar.text = "Primary Actor\nBudgetWise User"
    set_run_style(ar, size=11, bold=True, color=TEXT_DARK)

    uc1 = add_use_case_oval(slide, 3.35, 2.0, 2.2, 1.0, "View Dashboard")
    uc2 = add_use_case_oval(slide, 5.95, 1.9, 2.4, 1.0, "Manage Budgets")
    uc3 = add_use_case_oval(slide, 8.85, 2.0, 2.8, 1.0, "Track Expenses\n(Add/Edit/Delete/Filter)")
    uc4 = add_use_case_oval(slide, 4.15, 4.0, 2.4, 1.0, "View Reports")
    uc5 = add_use_case_oval(slide, 7.05, 3.9, 2.9, 1.0, "Receive Alerts and\nRecommendations")

    actor_center_x = Inches(2.4)
    actor_center_y = Inches(3.3)
    targets = [
        (uc1.left + uc1.width // 2, uc1.top + uc1.height // 2),
        (uc2.left + uc2.width // 2, uc2.top + uc2.height // 2),
        (uc3.left + uc3.width // 2, uc3.top + uc3.height // 2),
        (uc4.left + uc4.width // 2, uc4.top + uc4.height // 2),
        (uc5.left + uc5.width // 2, uc5.top + uc5.height // 2),
    ]
    for tx, ty in targets:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, actor_center_x, actor_center_y, tx, ty)
        line.line.color.rgb = TEXT_MUTED
        line.line.width = Pt(1.1)


def add_entity(slide, x, y, w, h, title, attrs, header_color):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box.line.color.rgb = BORDER
    box.line.width = Pt(1.3)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.52))
    header.fill.solid()
    header.fill.fore_color.rgb = header_color
    header.line.fill.background()

    ht = slide.shapes.add_textbox(Inches(x + 0.08), Inches(y + 0.1), Inches(w - 0.16), Inches(0.3))
    hp = ht.text_frame.paragraphs[0]
    hr = hp.add_run()
    hr.text = title
    set_run_style(hr, size=13, bold=True, color=RGBColor(255, 255, 255))

    content = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.65), Inches(w - 0.3), Inches(h - 0.8))
    tf = content.text_frame
    tf.clear()
    for i, attr in enumerate(attrs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = attr
        set_run_style(r, size=10.5, color=TEXT_DARK)
        p.space_after = Pt(3)


def create_er_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    add_header(
        slide,
        "ER Diagram",
        "Data entities inferred from Mongoose schemas in backend/models",
    )

    add_entity(
        slide,
        1.15,
        1.75,
        4.7,
        4.6,
        "Budget",
        [
            "PK: _id",
            "category (unique, enum)",
            "label",
            "icon",
            "color",
            "monthlyLimit",
            "createdAt, updatedAt",
        ],
        BLUE,
    )

    add_entity(
        slide,
        7.45,
        1.75,
        4.7,
        4.6,
        "Expense",
        [
            "PK: _id",
            "category (enum, maps to Budget.category)",
            "description",
            "amount",
            "date",
            "notes",
            "createdAt, updatedAt",
            "Indexes: date, category+date",
        ],
        GREEN,
    )

    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(5.85),
        Inches(3.95),
        Inches(7.45),
        Inches(3.95),
    )
    connector.line.color.rgb = TEXT_DARK
    connector.line.width = Pt(1.8)

    cardinality1 = slide.shapes.add_textbox(Inches(5.55), Inches(3.62), Inches(0.3), Inches(0.3))
    cp1 = cardinality1.text_frame.paragraphs[0]
    cp1.alignment = PP_ALIGN.CENTER
    cr1 = cp1.add_run()
    cr1.text = "1"
    set_run_style(cr1, size=12, bold=True, color=TEXT_DARK)

    cardinalityn = slide.shapes.add_textbox(Inches(7.5), Inches(3.62), Inches(0.4), Inches(0.3))
    cpn = cardinalityn.text_frame.paragraphs[0]
    cpn.alignment = PP_ALIGN.CENTER
    crn = cpn.add_run()
    crn.text = "N"
    set_run_style(crn, size=12, bold=True, color=TEXT_DARK)

    rel_label = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.85), Inches(4.15), Inches(1.55), Inches(0.58))
    rel_label.fill.solid()
    rel_label.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rel_label.line.color.rgb = BORDER
    rel_label.line.width = Pt(1)
    rtf = rel_label.text_frame
    rtf.clear()
    rp = rtf.paragraphs[0]
    rp.alignment = PP_ALIGN.CENTER
    rr = rp.add_run()
    rr.text = "category link"
    set_run_style(rr, size=10, bold=True, color=TEXT_MUTED)

    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.15), Inches(6.55), Inches(10.95), Inches(0.7))
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(255, 255, 255)
    note.line.color.rgb = BORDER
    note.line.width = Pt(1)
    ntf = note.text_frame
    ntf.clear()
    np = ntf.paragraphs[0]
    nr = np.add_run()
    nr.text = "Relationship meaning: one budget category can have many expenses in that category."
    set_run_style(nr, size=11, color=TEXT_DARK)


def create_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    add_header(
        slide,
        "Summary",
        "Project structure and behavior captured directly from implementation",
    )

    add_card(
        slide,
        0.9,
        1.75,
        11.55,
        4.85,
        "What this presentation documents",
        "- Layered architecture from React UI to MongoDB storage\n"
        "- Six modules aligned to backend routes and frontend pages\n"
        "- Operational data flow for create/read/update/delete and analytics\n"
        "- Primary user interactions as use cases\n"
        "- ER model based on Mongoose schemas\n\n"
        "Source references reviewed:\n"
        "backend/server.js, backend/routes/*.js, backend/models/*.js,\n"
        "frontend/src/pages/*.jsx, frontend/src/utils/api.js",
        fill=RGBColor(255, 255, 255),
    )


def build_presentation(output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs)
    create_architecture_slide(prs)
    create_modules_slide(prs)
    create_data_flow_slide(prs)
    create_use_case_slide(prs)
    create_er_slide(prs)
    create_summary_slide(prs)

    for i, slide in enumerate(prs.slides, start=1):
        add_slide_number(slide, i)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


if __name__ == "__main__":
    root = Path(__file__).resolve().parents[1]
    output_file = root / "BudgetWise_Project_Overview.pptx"
    build_presentation(output_file)
    print(f"Created presentation: {output_file}")
