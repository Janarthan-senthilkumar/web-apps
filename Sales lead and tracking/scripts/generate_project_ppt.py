from datetime import date

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# Slide constants (16:9 layout)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Color palette
BG = RGBColor(245, 248, 252)
PRIMARY = RGBColor(15, 44, 84)
SECONDARY = RGBColor(46, 134, 193)
ACCENT = RGBColor(25, 111, 61)
TEXT_DARK = RGBColor(35, 39, 42)
TEXT_MUTED = RGBColor(93, 109, 126)
CARD_BG = RGBColor(255, 255, 255)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def style_textframe(tf, font_size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    for paragraph in tf.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color


def set_paragraph(paragraph, text, size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    paragraph.text = text
    paragraph.alignment = align
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_title_block(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.2), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    set_paragraph(p, title, size=30, bold=True, color=PRIMARY)

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(12.2), Inches(0.6))
        stf = subtitle_box.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        set_paragraph(sp, subtitle, size=14, bold=False, color=TEXT_MUTED)

    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.55), Inches(3.5), Inches(0.07))
    line.fill.solid()
    line.fill.fore_color.rgb = SECONDARY
    line.line.fill.background()


def add_card(slide, left, top, width, height, title, body_lines, border=SECONDARY):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = CARD_BG
    rect.line.color.rgb = border
    rect.line.width = Pt(1.5)

    tf = rect.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    set_paragraph(p0, title, size=14, bold=True, color=PRIMARY)

    for line in body_lines:
        p = tf.add_paragraph()
        set_paragraph(p, line, size=11, color=TEXT_DARK)
        p.level = 0

    return rect


def add_arrow(slide, left, top, width, height):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = SECONDARY
    arrow.line.color.rgb = SECONDARY
    return arrow


def add_v_arrow(slide, left, top, width, height):
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = SECONDARY
    arrow.line.color.rgb = SECONDARY
    return arrow


def add_line(slide, x1, y1, x2, y2, color=TEXT_MUTED, width=1.5):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def build_presentation(output_file):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]

    # Slide 1: Cover
    s1 = prs.slides.add_slide(blank)
    set_slide_bg(s1, BG)

    hero = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.6))
    hero.fill.solid()
    hero.fill.fore_color.rgb = PRIMARY
    hero.line.fill.background()

    tbox = s1.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(11.5), Inches(0.8))
    tf = tbox.text_frame
    tf.clear()
    set_paragraph(tf.paragraphs[0], "Sales Lead Tracking & CRM Management System", size=34, bold=True, color=RGBColor(255, 255, 255))

    sub = s1.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(10.5), Inches(0.45))
    stf = sub.text_frame
    stf.clear()
    set_paragraph(stf.paragraphs[0], "Architecture | Modules | Data Flow | Use Cases | ER Diagram", size=14, color=RGBColor(220, 232, 247))

    details = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.3), Inches(11.7), Inches(3.9))
    details.fill.solid()
    details.fill.fore_color.rgb = CARD_BG
    details.line.color.rgb = SECONDARY
    details.line.width = Pt(2)

    dtf = details.text_frame
    dtf.clear()
    set_paragraph(dtf.paragraphs[0], "Project Snapshot", size=20, bold=True, color=PRIMARY)
    for line in [
        "Frontend: React (Vite), React Router, Axios, AuthContext",
        "Backend: Node.js, Express, JWT authentication, role-based authorization",
        "Database: MongoDB with Mongoose models (User, Lead, Customer, FollowUp, Activity)",
        "Purpose: Manage lead lifecycle from capture to conversion and post-sale relationship",
        f"Prepared on: {date.today().isoformat()}",
    ]:
        p = dtf.add_paragraph()
        set_paragraph(p, f"- {line}", size=14, color=TEXT_DARK)

    # Slide 2: Architecture Diagram
    s2 = prs.slides.add_slide(blank)
    set_slide_bg(s2, BG)
    add_title_block(s2, "1) System Architecture Diagram", "MERN-based role-aware CRM architecture")

    users = add_card(
        s2,
        Inches(0.5),
        Inches(2.0),
        Inches(2.1),
        Inches(1.4),
        "User Roles",
        ["Admin", "Manager", "Executive"],
        border=PRIMARY,
    )
    frontend = add_card(
        s2,
        Inches(3.0),
        Inches(1.85),
        Inches(3.1),
        Inches(1.8),
        "React Frontend",
        ["Pages: Dashboard, Leads, Customers,", "FollowUps, Team, Login", "AuthContext + ProtectedRoute", "Axios API Client + JWT"],
        border=SECONDARY,
    )
    api = add_card(
        s2,
        Inches(6.5),
        Inches(1.85),
        Inches(2.8),
        Inches(1.8),
        "Express API",
        ["Routes + Controllers", "Middleware: auth, errorHandler", "Service: activityService"],
        border=SECONDARY,
    )
    db = add_card(
        s2,
        Inches(9.8),
        Inches(1.85),
        Inches(3.0),
        Inches(1.8),
        "MongoDB",
        ["Collections:", "users, leads, customers,", "followups, activities"],
        border=ACCENT,
    )

    add_arrow(s2, Inches(2.65), Inches(2.35), Inches(0.28), Inches(0.35))
    add_arrow(s2, Inches(6.15), Inches(2.35), Inches(0.28), Inches(0.35))
    add_arrow(s2, Inches(9.45), Inches(2.35), Inches(0.28), Inches(0.35))

    mod_title = s2.shapes.add_textbox(Inches(0.6), Inches(4.05), Inches(12), Inches(0.35))
    set_paragraph(mod_title.text_frame.paragraphs[0], "Backend Domain Modules", size=16, bold=True, color=PRIMARY)

    mod_w = Inches(2.05)
    mod_h = Inches(1.15)
    mod_top = Inches(4.4)
    mod_lefts = [Inches(0.6), Inches(2.85), Inches(5.1), Inches(7.35), Inches(9.6)]
    modules = [
        ("Auth", ["login/register", "JWT + roles"]),
        ("Lead Mgmt", ["capture, assign", "status, convert"]),
        ("Customer Mgmt", ["customer records", "interactions"]),
        ("Follow-Up", ["schedule, complete", "queue tracking"]),
        ("Dashboard", ["summary KPIs", "conversion report"]),
    ]

    for idx, (title, lines) in enumerate(modules):
        add_card(s2, mod_lefts[idx], mod_top, mod_w, mod_h, title, lines, border=SECONDARY)

    # Slide 3: Six Modules Breakdown
    s3 = prs.slides.add_slide(blank)
    set_slide_bg(s3, BG)
    add_title_block(s3, "2) Project Breakdown into 6 Modules", "Functional decomposition aligned to client pages and server controllers")

    module_cards = [
        ("Module 1: Authentication", ["Register/Login/Bootstrap Admin", "Token management + /auth/me sync", "JWT verification middleware"]),
        ("Module 2: Team & RBAC", ["Roles: admin, manager, executive", "Create/list users in Team page", "Route-level role authorization"]),
        ("Module 3: Lead Management", ["Create/update/search leads", "Assign ownership + update status", "Convert WON leads to customers"]),
        ("Module 4: Follow-Up Management", ["Schedule follow-ups by lead", "Track due dates and status", "Complete follow-up with outcomes"]),
        ("Module 5: Customer CRM", ["Create/update customers", "Track lifecycle stage", "Interaction history timeline"]),
        ("Module 6: Dashboard & Analytics", ["Summary KPIs and pipeline value", "Monthly conversion reporting", "Executive performance + activity feed"]),
    ]

    x_positions = [Inches(0.55), Inches(6.85)]
    y_start = Inches(1.9)
    card_w = Inches(5.95)
    card_h = Inches(1.55)
    gap = Inches(0.35)

    for i, (title, lines) in enumerate(module_cards):
        col = i % 2
        row = i // 2
        x = x_positions[col]
        y = y_start + row * (card_h + gap)
        border = PRIMARY if i % 3 == 0 else SECONDARY
        add_card(s3, x, y, card_w, card_h, title, lines, border=border)

    # Slide 4: Module Explanation Table
    s4 = prs.slides.add_slide(blank)
    set_slide_bg(s4, BG)
    add_title_block(s4, "Module Explanations (Responsibilities and Key Files)")

    table = s4.shapes.add_table(7, 4, Inches(0.5), Inches(1.8), Inches(12.3), Inches(5.2)).table
    table.columns[0].width = Inches(2.25)
    table.columns[1].width = Inches(3.15)
    table.columns[2].width = Inches(3.2)
    table.columns[3].width = Inches(3.7)

    headers = ["Module", "Primary Responsibility", "Key Backend Paths", "Key Frontend Paths"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(12)
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.name = "Segoe UI"
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY

    rows = [
        ["Authentication", "Session creation and identity checks", "authController, authRoutes, middleware/auth", "context/AuthContext, pages/LoginPage"],
        ["Team & RBAC", "Role-based access and team management", "authRoutes (/users), authorize middleware", "components/ProtectedRoute, pages/TeamPage"],
        ["Lead Mgmt", "Lead lifecycle and conversion", "leadController, models/Lead", "pages/LeadsPage"],
        ["Follow-Up", "Scheduling and completion workflow", "followUpController, models/FollowUp", "pages/FollowUpsPage"],
        ["Customer CRM", "Customer records and interactions", "customerController, models/Customer", "pages/CustomersPage"],
        ["Dashboard", "Metrics and reporting aggregates", "dashboardController, models/Activity", "pages/DashboardPage"],
    ]

    for r, row_data in enumerate(rows, start=1):
        for c, value in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = value
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10.5)
                    run.font.name = "Segoe UI"
                    run.font.color.rgb = TEXT_DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(250, 252, 255) if r % 2 == 1 else RGBColor(240, 246, 255)

    # Slide 5: Data Flow
    s5 = prs.slides.add_slide(blank)
    set_slide_bg(s5, BG)
    add_title_block(s5, "3) Data Flow", "Lead-to-Customer operational flow implemented by routes and controllers")

    steps = [
        ("1. Authenticate", "User logs in\nJWT stored in localStorage"),
        ("2. Create Lead", "POST /leads\nLead assigned + activity logged"),
        ("3. Work Lead", "PATCH status/assign\nSearch + ownership filtering"),
        ("4. Schedule Follow-Up", "POST /followups\nLead.nextFollowUpDate updated"),
        ("5. Complete Follow-Up", "PATCH /followups/:id/complete\nLead.lastFollowUpAt updated"),
        ("6. Convert & Report", "POST /leads/:id/convert\nCustomer created + dashboard aggregates"),
    ]

    start_x = Inches(0.55)
    top_y = Inches(2.2)
    box_w = Inches(2.0)
    box_h = Inches(1.4)
    gap_x = Inches(0.15)

    for i, (title, text) in enumerate(steps):
        x = start_x + i * (box_w + gap_x)
        box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top_y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.color.rgb = SECONDARY
        box.line.width = Pt(1.4)
        tf = box.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        set_paragraph(p0, title, size=11.5, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
        p1 = tf.add_paragraph()
        set_paragraph(p1, text, size=9.5, color=TEXT_DARK, align=PP_ALIGN.CENTER)

        if i < len(steps) - 1:
            add_arrow(
                s5,
                x + box_w + Inches(0.03),
                top_y + Inches(0.52),
                Inches(0.1),
                Inches(0.33),
            )

    notes = add_card(
        s5,
        Inches(0.9),
        Inches(4.35),
        Inches(11.7),
        Inches(2.2),
        "How state propagates through the system",
        [
            "Frontend sends authenticated API requests using Axios interceptor with Bearer token.",
            "Controllers validate role permissions, mutate MongoDB entities, and write activity logs.",
            "Follow-up completion updates lead timing fields; conversion creates linked customer records.",
            "Dashboard endpoints aggregate Leads + Customers + FollowUps + Activities for KPIs.",
        ],
        border=ACCENT,
    )
    notes.text_frame.word_wrap = True

    # Slide 6: Use Case
    s6 = prs.slides.add_slide(blank)
    set_slide_bg(s6, BG)
    add_title_block(s6, "4) Use Case Diagram", "Actors and major business operations")

    # Actors
    actor_admin = add_card(s6, Inches(0.45), Inches(2.0), Inches(1.9), Inches(1.05), "Admin", ["Full control"], border=PRIMARY)
    actor_mgr = add_card(s6, Inches(0.45), Inches(3.35), Inches(1.9), Inches(1.05), "Manager", ["Operational control"], border=SECONDARY)
    actor_exec = add_card(s6, Inches(0.45), Inches(4.7), Inches(1.9), Inches(1.05), "Executive", ["Execution role"], border=ACCENT)

    system = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.8), Inches(1.75), Inches(9.9), Inches(4.4))
    system.fill.solid()
    system.fill.fore_color.rgb = RGBColor(255, 255, 255)
    system.line.color.rgb = PRIMARY
    system.line.width = Pt(2)

    sys_title = s6.shapes.add_textbox(Inches(3.0), Inches(1.9), Inches(4.5), Inches(0.35))
    set_paragraph(sys_title.text_frame.paragraphs[0], "LeadFlow CRM System", size=15, bold=True, color=PRIMARY)

    uc_positions = [
        (Inches(3.1), Inches(2.35), "Login / Register"),
        (Inches(5.4), Inches(2.35), "Manage Team"),
        (Inches(8.0), Inches(2.35), "Create / Assign Leads"),
        (Inches(3.1), Inches(3.45), "Update Lead Status"),
        (Inches(5.6), Inches(3.45), "Schedule / Complete Follow-Up"),
        (Inches(8.35), Inches(3.45), "Convert Lead to Customer"),
        (Inches(4.2), Inches(4.55), "Manage Customers & Interactions"),
        (Inches(7.4), Inches(4.55), "View Dashboard Reports"),
    ]

    uc_shapes = []
    for x, y, text in uc_positions:
        oval = s6.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(2.25), Inches(0.72))
        oval.fill.solid()
        oval.fill.fore_color.rgb = RGBColor(234, 243, 252)
        oval.line.color.rgb = SECONDARY
        tf = oval.text_frame
        tf.clear()
        set_paragraph(tf.paragraphs[0], text, size=10.3, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
        uc_shapes.append(oval)

    # Actor connectors
    add_line(s6, Inches(2.35), Inches(2.53), Inches(3.1), Inches(2.72), color=PRIMARY)
    add_line(s6, Inches(2.35), Inches(2.53), Inches(5.4), Inches(2.72), color=PRIMARY)
    add_line(s6, Inches(2.35), Inches(2.53), Inches(7.95), Inches(2.72), color=PRIMARY)
    add_line(s6, Inches(2.35), Inches(3.88), Inches(3.1), Inches(3.83), color=SECONDARY)
    add_line(s6, Inches(2.35), Inches(3.88), Inches(5.6), Inches(3.83), color=SECONDARY)
    add_line(s6, Inches(2.35), Inches(3.88), Inches(7.35), Inches(4.9), color=SECONDARY)
    add_line(s6, Inches(2.35), Inches(5.23), Inches(3.1), Inches(3.83), color=ACCENT)
    add_line(s6, Inches(2.35), Inches(5.23), Inches(5.6), Inches(3.83), color=ACCENT)
    add_line(s6, Inches(2.35), Inches(5.23), Inches(4.2), Inches(4.9), color=ACCENT)

    # Slide 7: ER Diagram
    s7 = prs.slides.add_slide(blank)
    set_slide_bg(s7, BG)
    add_title_block(s7, "5) ER Diagram", "MongoDB entity relationship model from Mongoose schemas")

    # Entity boxes
    user_box = add_card(
        s7,
        Inches(0.35),
        Inches(2.0),
        Inches(2.55),
        Inches(2.2),
        "User",
        ["PK _id", "name, email, password", "role, isActive", "createdAt, updatedAt"],
        border=PRIMARY,
    )
    lead_box = add_card(
        s7,
        Inches(3.2),
        Inches(1.75),
        Inches(3.0),
        Inches(2.7),
        "Lead",
        [
            "PK _id",
            "fullName, email, phone, company",
            "status, source, estimatedValue",
            "FK assignedTo -> User",
            "FK createdBy -> User",
            "FK convertedCustomer -> Customer",
        ],
        border=SECONDARY,
    )
    followup_box = add_card(
        s7,
        Inches(6.55),
        Inches(1.75),
        Inches(2.8),
        Inches(2.6),
        "FollowUp",
        ["PK _id", "FK lead -> Lead", "FK assignedTo -> User", "dueDate, type, status", "FK createdBy -> User"],
        border=SECONDARY,
    )
    customer_box = add_card(
        s7,
        Inches(9.7),
        Inches(1.75),
        Inches(3.2),
        Inches(2.8),
        "Customer",
        [
            "PK _id",
            "fullName, company, stage, value",
            "FK leadRef -> Lead",
            "FK accountManager -> User",
            "interactions[] {type, summary, createdBy}",
        ],
        border=ACCENT,
    )
    activity_box = add_card(
        s7,
        Inches(4.5),
        Inches(4.75),
        Inches(4.25),
        Inches(1.9),
        "Activity",
        ["PK _id", "action, entityType, entityId", "FK actor -> User", "metadata{} , createdAt"],
        border=RGBColor(120, 81, 169),
    )

    # Relationship lines + labels
    add_line(s7, Inches(2.9), Inches(2.6), Inches(3.2), Inches(2.6), color=PRIMARY, width=2)
    add_line(s7, Inches(2.9), Inches(3.2), Inches(3.2), Inches(3.2), color=PRIMARY, width=2)
    add_line(s7, Inches(6.2), Inches(2.8), Inches(6.55), Inches(2.8), color=SECONDARY, width=2)
    add_line(s7, Inches(9.35), Inches(2.8), Inches(9.7), Inches(2.8), color=SECONDARY, width=2)
    add_line(s7, Inches(1.65), Inches(4.2), Inches(5.2), Inches(4.75), color=TEXT_MUTED, width=1.8)
    add_line(s7, Inches(7.95), Inches(4.35), Inches(8.35), Inches(4.75), color=TEXT_MUTED, width=1.8)

    label1 = s7.shapes.add_textbox(Inches(2.95), Inches(2.35), Inches(1.0), Inches(0.25))
    set_paragraph(label1.text_frame.paragraphs[0], "1 : many", size=8.8, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    label2 = s7.shapes.add_textbox(Inches(6.2), Inches(2.55), Inches(1.1), Inches(0.25))
    set_paragraph(label2.text_frame.paragraphs[0], "Lead 1:M FollowUp", size=8.8, color=TEXT_MUTED)

    label3 = s7.shapes.add_textbox(Inches(9.15), Inches(2.55), Inches(1.5), Inches(0.25))
    set_paragraph(label3.text_frame.paragraphs[0], "Lead 1:0..1 Customer", size=8.8, color=TEXT_MUTED)

    label4 = s7.shapes.add_textbox(Inches(4.2), Inches(4.35), Inches(3.2), Inches(0.25))
    set_paragraph(label4.text_frame.paragraphs[0], "User 1:M Activity (actor)", size=8.8, color=TEXT_MUTED)

    label5 = s7.shapes.add_textbox(Inches(8.35), Inches(4.35), Inches(3.0), Inches(0.25))
    set_paragraph(label5.text_frame.paragraphs[0], "Activity references entityId", size=8.8, color=TEXT_MUTED)

    # Slide 8: Closing summary
    s8 = prs.slides.add_slide(blank)
    set_slide_bg(s8, BG)
    add_title_block(s8, "Solution Summary")

    add_card(
        s8,
        Inches(0.75),
        Inches(2.0),
        Inches(12.0),
        Inches(3.5),
        "What this architecture enables",
        [
            "Structured lead lifecycle: capture -> assign -> follow-up -> convert -> customer management.",
            "Role-aware operations with secure JWT authentication and backend authorization guards.",
            "Traceability through centralized activity logging across auth, leads, customers, and follow-ups.",
            "Actionable business insights through dashboard KPIs, conversion analytics, and performance reports.",
        ],
        border=PRIMARY,
    )

    thanks = s8.shapes.add_textbox(Inches(0.75), Inches(6.1), Inches(12.0), Inches(0.6))
    tf = thanks.text_frame
    tf.clear()
    set_paragraph(tf.paragraphs[0], "End of Presentation", size=20, bold=True, color=SECONDARY, align=PP_ALIGN.CENTER)

    prs.save(output_file)


if __name__ == "__main__":
    build_presentation("Sales_Lead_Tracking_Project_Overview.pptx")
