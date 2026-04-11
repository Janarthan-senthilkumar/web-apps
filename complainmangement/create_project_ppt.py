from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_FILE = Path("Complaint_Management_Project_Overview.pptx")


# Theme colors
NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
ACCENT = RGBColor(37, 99, 235)
ACCENT_DARK = RGBColor(30, 64, 175)
GREEN = RGBColor(16, 185, 129)
AMBER = RGBColor(245, 158, 11)
PURPLE = RGBColor(124, 58, 237)
LIGHT_BG = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)


def set_text_frame(
    text_frame,
    text,
    font_size=16,
    bold=False,
    color=NAVY,
    align=PP_ALIGN.LEFT,
):
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title_block(slide, title, subtitle):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7))
    set_text_frame(
        title_box.text_frame,
        title,
        font_size=30,
        bold=True,
        color=NAVY,
        align=PP_ALIGN.LEFT,
    )
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.4))
    set_text_frame(
        sub_box.text_frame,
        subtitle,
        font_size=14,
        bold=False,
        color=MUTED,
        align=PP_ALIGN.LEFT,
    )


def add_card(slide, left, top, width, height, title, lines, fill_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(226, 232, 240)
    shape.line.width = Pt(1.2)

    title_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.12), width - Inches(0.2), Inches(0.4))
    set_text_frame(title_box.text_frame, title, font_size=15, bold=True, color=ACCENT_DARK)

    body = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.5), width - Inches(0.3), height - Inches(0.6))
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.bullet = True
        p.font.size = Pt(12.5)
        p.font.color.rgb = SLATE
        p.space_after = Pt(8)


def add_arrow(slide, left, top, width, height, text="", color=ACCENT):
    arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    if text:
        set_text_frame(
            arrow.text_frame,
            text,
            font_size=10,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )
    return arrow


def add_entity_box(slide, left, top, width, height, title, fields, title_color):
    outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    outer.fill.solid()
    outer.fill.fore_color.rgb = WHITE
    outer.line.color.rgb = RGBColor(203, 213, 225)
    outer.line.width = Pt(1.2)

    header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left + Inches(0.02),
        top + Inches(0.02),
        width - Inches(0.04),
        Inches(0.4),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = title_color
    header.line.fill.background()
    set_text_frame(header.text_frame, title, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    body = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.48), width - Inches(0.16), height - Inches(0.56))
    tf = body.text_frame
    tf.clear()
    for idx, field in enumerate(fields):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = field
        p.font.size = Pt(10.5)
        p.font.color.rgb = SLATE
        p.space_after = Pt(4)


def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()

    ribbon = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.45))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = ACCENT_DARK
    ribbon.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.8), Inches(1.2))
    set_text_frame(
        title.text_frame,
        "College Complaint Management System",
        font_size=42,
        bold=True,
        color=NAVY,
        align=PP_ALIGN.LEFT,
    )

    sub = slide.shapes.add_textbox(Inches(0.82), Inches(2.45), Inches(11.2), Inches(0.8))
    set_text_frame(
        sub.text_frame,
        "Project Architecture and Functional Design Overview",
        font_size=20,
        bold=False,
        color=SLATE,
        align=PP_ALIGN.LEFT,
    )

    details = slide.shapes.add_textbox(Inches(0.82), Inches(3.2), Inches(8.5), Inches(1.8))
    tf = details.text_frame
    tf.clear()
    lines = [
        "Technology: MERN stack (MongoDB, Express, React, Node.js)",
        "Roles: Admin, Staff, Student",
        "Includes architecture, modules, data flow, use case, and ER diagram",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.bullet = True
        p.font.size = Pt(16)
        p.font.color.rgb = SLATE
        p.space_after = Pt(12)

    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(5.7), Inches(4.9), Inches(0.6))
    pill.fill.solid()
    pill.fill.fore_color.rgb = ACCENT
    pill.line.fill.background()
    set_text_frame(
        pill.text_frame,
        "Prepared from project source analysis",
        font_size=13,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )


def create_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "1) Architecture Diagram", "End-to-end system view across client, API, realtime, storage, and database layers")

    # Actor and frontend
    add_card(
        slide,
        Inches(0.6),
        Inches(1.5),
        Inches(2.15),
        Inches(1.15),
        "Users",
        ["Admin", "Staff", "Student"],
        fill_color=RGBColor(239, 246, 255),
    )
    add_arrow(slide, Inches(2.85), Inches(1.84), Inches(0.65), Inches(0.45), "", color=ACCENT)
    add_card(
        slide,
        Inches(3.55),
        Inches(1.5),
        Inches(2.8),
        Inches(1.15),
        "Frontend Layer",
        ["React + Vite UI", "Role-based routes", "Axios API client"],
        fill_color=RGBColor(236, 253, 245),
    )

    add_arrow(slide, Inches(6.4), Inches(1.84), Inches(0.65), Inches(0.45), "HTTPS", color=ACCENT)
    add_card(
        slide,
        Inches(7.15),
        Inches(1.5),
        Inches(2.95),
        Inches(1.15),
        "Backend API Layer",
        ["Node.js + Express", "JWT auth middleware", "Route controllers"],
        fill_color=RGBColor(245, 243, 255),
    )

    # Subsystems
    add_card(
        slide,
        Inches(1.0),
        Inches(3.35),
        Inches(2.7),
        Inches(1.1),
        "Realtime Channel",
        ["Socket.io events", "newComplaint", "complaintUpdated/newResponse"],
        fill_color=RGBColor(254, 249, 195),
    )
    add_card(
        slide,
        Inches(4.1),
        Inches(3.35),
        Inches(2.7),
        Inches(1.1),
        "File Storage",
        ["Multer upload middleware", "3 files max", "/uploads static serving"],
        fill_color=RGBColor(255, 237, 213),
    )
    add_card(
        slide,
        Inches(7.2),
        Inches(3.35),
        Inches(2.7),
        Inches(1.1),
        "Data Layer",
        ["Mongoose models", "User, Complaint", "Response, Notification"],
        fill_color=RGBColor(224, 231, 255),
    )

    db = slide.shapes.add_shape(MSO_SHAPE.CAN, Inches(10.35), Inches(3.45), Inches(2.3), Inches(1.0))
    db.fill.solid()
    db.fill.fore_color.rgb = RGBColor(219, 234, 254)
    db.line.color.rgb = RGBColor(147, 197, 253)
    set_text_frame(db.text_frame, "MongoDB\ncomplaint_management", font_size=12, bold=True, color=ACCENT_DARK, align=PP_ALIGN.CENTER)

    # Connectors
    conn1 = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(8.6),
        Inches(2.65),
        Inches(8.55),
        Inches(3.35),
    )
    conn1.line.color.rgb = MUTED
    conn1.line.width = Pt(1.8)

    conn2 = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(8.0),
        Inches(2.65),
        Inches(5.45),
        Inches(3.35),
    )
    conn2.line.color.rgb = MUTED
    conn2.line.width = Pt(1.8)

    conn3 = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(7.55),
        Inches(2.65),
        Inches(2.35),
        Inches(3.35),
    )
    conn3.line.color.rgb = MUTED
    conn3.line.width = Pt(1.8)

    conn4 = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(9.9),
        Inches(3.9),
        Inches(10.35),
        Inches(3.9),
    )
    conn4.line.color.rgb = MUTED
    conn4.line.width = Pt(1.8)

    legend = slide.shapes.add_textbox(Inches(0.6), Inches(5.95), Inches(12.1), Inches(1.0))
    tf = legend.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = "Key idea: Frontend and backend are decoupled; backend centralizes auth, business rules, notifications, and persistence."
    p1.font.size = Pt(12.5)
    p1.font.color.rgb = SLATE
    p1.space_after = Pt(6)
    p2 = tf.add_paragraph()
    p2.text = "Realtime events (Socket.io) and periodic polling are both used to keep dashboards current."
    p2.font.size = Pt(12.5)
    p2.font.color.rgb = SLATE


def create_module_slides(prs):
    # Slide 1: Modules 1-3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "2) Project Broken Into 6 Modules", "Modules 1-3: Core identity, complaint intake, and admin orchestration")

    add_card(
        slide,
        Inches(0.6),
        Inches(1.4),
        Inches(4.0),
        Inches(1.8),
        "Module 1: Authentication & Access Control",
        [
            "Registration/login with JWT token issuance",
            "Role enforcement using protect + authorize middleware",
            "Frontend role routes for admin/staff/user pages",
        ],
    )
    add_card(
        slide,
        Inches(4.9),
        Inches(1.4),
        Inches(4.0),
        Inches(1.8),
        "Module 2: Complaint Intake & Tracking",
        [
            "Student submits complaint with category, priority, attachments",
            "Multer validates file type/size and stores in uploads/",
            "Students fetch only their own complaints with timeline view",
        ],
    )
    add_card(
        slide,
        Inches(9.2),
        Inches(1.4),
        Inches(3.5),
        Inches(1.8),
        "Module 3: Admin Operations",
        [
            "View/filter/paginate all complaints",
            "Assign staff and change complaint status",
            "Manage users: edit role, department, active state, delete",
        ],
    )

    files = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(3.55), Inches(12.1), Inches(2.7))
    files.fill.solid()
    files.fill.fore_color.rgb = RGBColor(241, 245, 249)
    files.line.color.rgb = RGBColor(203, 213, 225)
    files.line.width = Pt(1.0)

    files_title = slide.shapes.add_textbox(Inches(0.8), Inches(3.72), Inches(11.7), Inches(0.4))
    set_text_frame(files_title.text_frame, "Primary backend/frontend files involved", font_size=14, bold=True, color=ACCENT_DARK)

    file_list = slide.shapes.add_textbox(Inches(0.8), Inches(4.1), Inches(11.7), Inches(2.0))
    tf = file_list.text_frame
    tf.clear()
    lines = [
        "Module 1: routes/auth.js, middleware/auth.js, models/User.js, context/AuthContext.jsx, components/ProtectedRoute.jsx",
        "Module 2: routes/complaints.js, middleware/upload.js, models/Complaint.js, pages/user/SubmitComplaint.jsx, pages/user/TrackComplaints.jsx",
        "Module 3: routes/admin.js, pages/admin/AdminComplaints.jsx, pages/admin/AdminUsers.jsx, pages/admin/AdminStaff.jsx",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.bullet = True
        p.font.size = Pt(12)
        p.font.color.rgb = SLATE
        p.space_after = Pt(8)

    # Slide 2: Modules 4-6
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide2, "2) Project Broken Into 6 Modules", "Modules 4-6: Staff handling, communication loop, and reporting insights")

    add_card(
        slide2,
        Inches(0.6),
        Inches(1.4),
        Inches(4.0),
        Inches(1.8),
        "Module 4: Staff Resolution Workspace",
        [
            "Staff sees assigned complaints only",
            "Updates status and internal notes",
            "Role-safe actions enforced at API level",
        ],
    )
    add_card(
        slide2,
        Inches(4.9),
        Inches(1.4),
        Inches(4.0),
        Inches(1.8),
        "Module 5: Responses & Notifications",
        [
            "Threaded responses between user/staff/admin",
            "Internal notes hidden from student role",
            "Unread notification center with mark read/read-all",
        ],
    )
    add_card(
        slide2,
        Inches(9.2),
        Inches(1.4),
        Inches(3.5),
        Inches(1.8),
        "Module 6: Dashboards & Analytics",
        [
            "Admin stats via aggregation (status/category/priority)",
            "Staff and user dashboard KPIs",
            "Charts built using Recharts components",
        ],
    )

    files2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(3.55), Inches(12.1), Inches(2.7))
    files2.fill.solid()
    files2.fill.fore_color.rgb = RGBColor(241, 245, 249)
    files2.line.color.rgb = RGBColor(203, 213, 225)
    files2.line.width = Pt(1.0)

    files_title2 = slide2.shapes.add_textbox(Inches(0.8), Inches(3.72), Inches(11.7), Inches(0.4))
    set_text_frame(files_title2.text_frame, "Primary backend/frontend files involved", font_size=14, bold=True, color=ACCENT_DARK)

    file_list2 = slide2.shapes.add_textbox(Inches(0.8), Inches(4.1), Inches(11.7), Inches(2.0))
    tf2 = file_list2.text_frame
    tf2.clear()
    lines2 = [
        "Module 4: routes/staff.js, pages/staff/StaffDashboard.jsx, pages/staff/StaffComplaints.jsx",
        "Module 5: routes/responses.js, routes/notifications.js, models/Response.js, models/Notification.js, components/Topbar.jsx",
        "Module 6: routes/admin.js (/stats), pages/admin/AdminDashboard.jsx, pages/user/UserDashboard.jsx",
    ]
    for i, line in enumerate(lines2):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = line
        p.bullet = True
        p.font.size = Pt(12)
        p.font.color.rgb = SLATE
        p.space_after = Pt(8)


def create_data_flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "3) Data Flow", "Complaint lifecycle from submission to resolution and user feedback")

    steps = [
        ("Step 1", "Student submits complaint\n(+ optional attachments)", RGBColor(219, 234, 254)),
        ("Step 2", "API validates token,\nparses multipart files", RGBColor(220, 252, 231)),
        ("Step 3", "Complaint document saved\nin MongoDB", RGBColor(237, 233, 254)),
        ("Step 4", "Admin gets notification\nand assigns staff", RGBColor(254, 243, 199)),
        ("Step 5", "Staff updates status\nand posts response", RGBColor(255, 237, 213)),
        ("Step 6", "Notification + realtime event\nsent to related users", RGBColor(224, 231, 255)),
        ("Step 7", "Student tracks status,\nreads response thread", RGBColor(209, 250, 229)),
    ]

    y = Inches(1.6)
    for i, (label, text, color) in enumerate(steps):
        x = Inches(0.7 + i * 1.78)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.55), Inches(2.55))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = RGBColor(203, 213, 225)
        box.line.width = Pt(1.0)

        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.12), y + Inches(0.12), Inches(1.3), Inches(0.35))
        tag.fill.solid()
        tag.fill.fore_color.rgb = ACCENT_DARK
        tag.line.fill.background()
        set_text_frame(tag.text_frame, label, font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        desc = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.56), Inches(1.3), Inches(1.85))
        tf = desc.text_frame
        tf.clear()
        for idx, part in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = part
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(11)
            p.font.color.rgb = NAVY
            p.space_after = Pt(6)

        if i < len(steps) - 1:
            add_arrow(slide, x + Inches(1.58), y + Inches(1.05), Inches(0.18), Inches(0.45), "", color=ACCENT)

    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(4.45), Inches(11.9), Inches(1.85))
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(248, 250, 252)
    note.line.color.rgb = RGBColor(203, 213, 225)
    note.line.width = Pt(1.0)

    note_title = slide.shapes.add_textbox(Inches(0.92), Inches(4.62), Inches(11.5), Inches(0.35))
    set_text_frame(note_title.text_frame, "Data Handling Notes", font_size=14, bold=True, color=ACCENT_DARK)

    note_body = slide.shapes.add_textbox(Inches(0.92), Inches(4.95), Inches(11.5), Inches(1.2))
    tfb = note_body.text_frame
    tfb.clear()
    bullet_points = [
        "Authorization is enforced before sensitive read/write operations (JWT + role check).",
        "Internal staff notes are hidden from students when loading response threads.",
        "UI freshness combines polling (10-20 seconds) with Socket.io emits for key events.",
    ]
    for i, b in enumerate(bullet_points):
        p = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
        p.text = b
        p.bullet = True
        p.font.size = Pt(12)
        p.font.color.rgb = SLATE
        p.space_after = Pt(8)


def create_use_case_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "4) Use Case", "Role-wise functional capabilities in the complaint management workflow")

    # System boundary
    boundary = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.4), Inches(1.35), Inches(8.8), Inches(5.7))
    boundary.fill.solid()
    boundary.fill.fore_color.rgb = RGBColor(250, 250, 252)
    boundary.line.color.rgb = RGBColor(148, 163, 184)
    boundary.line.width = Pt(1.4)
    boundary_txt = slide.shapes.add_textbox(Inches(2.55), Inches(1.45), Inches(3.8), Inches(0.3))
    set_text_frame(boundary_txt.text_frame, "System Boundary: Complaint Management System", font_size=11, bold=True, color=MUTED)

    # Actors
    actors = [
        ("Student", Inches(0.55), Inches(2.0), RGBColor(16, 185, 129)),
        ("Staff", Inches(0.55), Inches(3.25), RGBColor(59, 130, 246)),
        ("Admin", Inches(0.55), Inches(4.5), RGBColor(124, 58, 237)),
    ]
    actor_positions = {}
    for name, x, y, color in actors:
        head = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.25), y, Inches(0.35), Inches(0.35))
        head.fill.solid()
        head.fill.fore_color.rgb = color
        head.line.fill.background()

        body = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            x + Inches(0.43),
            y + Inches(0.35),
            x + Inches(0.43),
            y + Inches(0.95),
        )
        body.line.color.rgb = color
        body.line.width = Pt(2.0)

        arms = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            x + Inches(0.15),
            y + Inches(0.55),
            x + Inches(0.71),
            y + Inches(0.55),
        )
        arms.line.color.rgb = color
        arms.line.width = Pt(2.0)

        leg1 = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            x + Inches(0.43),
            y + Inches(0.95),
            x + Inches(0.2),
            y + Inches(1.3),
        )
        leg1.line.color.rgb = color
        leg1.line.width = Pt(2.0)

        leg2 = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            x + Inches(0.43),
            y + Inches(0.95),
            x + Inches(0.66),
            y + Inches(1.3),
        )
        leg2.line.color.rgb = color
        leg2.line.width = Pt(2.0)

        label = slide.shapes.add_textbox(x, y + Inches(1.32), Inches(0.9), Inches(0.3))
        set_text_frame(label.text_frame, name, font_size=11, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
        actor_positions[name] = (x + Inches(0.72), y + Inches(0.55))

    # Use case ovals
    use_cases = [
        ("Authenticate\n(Login/Register)", Inches(3.0), Inches(1.95), Inches(2.2), Inches(0.9), ["Student", "Staff", "Admin"]),
        ("Submit\nComplaint", Inches(3.0), Inches(3.05), Inches(2.2), Inches(0.9), ["Student"]),
        ("Track Status &\nView Responses", Inches(3.0), Inches(4.15), Inches(2.2), Inches(0.9), ["Student"]),
        ("Handle Assigned\nComplaints", Inches(5.65), Inches(2.55), Inches(2.2), Inches(0.9), ["Staff"]),
        ("Update Status /\nAdd Notes", Inches(5.65), Inches(3.7), Inches(2.2), Inches(0.9), ["Staff"]),
        ("Assign Staff /\nManage Users", Inches(8.3), Inches(2.25), Inches(2.2), Inches(0.9), ["Admin"]),
        ("View Analytics &\nControl Workflow", Inches(8.3), Inches(3.65), Inches(2.2), Inches(0.9), ["Admin"]),
    ]

    oval_centers = []
    for text, x, y, w, h, actors_for_case in use_cases:
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
        oval.fill.solid()
        oval.fill.fore_color.rgb = RGBColor(226, 232, 240)
        oval.line.color.rgb = RGBColor(100, 116, 139)
        set_text_frame(oval.text_frame, text, font_size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

        cx = x + w / 2
        cy = y + h / 2
        oval_centers.append((actors_for_case, cx, cy))

    # Connections
    for actor_names, cx, cy in oval_centers:
        for actor_name in actor_names:
            ax, ay = actor_positions[actor_name]
            link = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, ax, ay, cx, cy)
            link.line.color.rgb = RGBColor(148, 163, 184)
            link.line.width = Pt(1.1)


def create_er_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "5) ER Diagram", "Core entities and relationships implemented in MongoDB via Mongoose schemas")

    add_entity_box(
        slide,
        Inches(0.7),
        Inches(1.55),
        Inches(3.2),
        Inches(2.8),
        "User",
        [
            "PK _id",
            "name",
            "email (unique)",
            "password (hashed)",
            "role {admin,staff,user}",
            "department",
            "isActive",
            "createdAt, updatedAt",
        ],
        title_color=PURPLE,
    )

    add_entity_box(
        slide,
        Inches(4.8),
        Inches(1.25),
        Inches(3.7),
        Inches(3.4),
        "Complaint",
        [
            "PK _id",
            "title, description",
            "category, priority, status",
            "attachments[]",
            "FK submittedBy -> User._id",
            "FK assignedTo -> User._id (nullable)",
            "department, notes, resolvedAt",
            "createdAt, updatedAt",
        ],
        title_color=ACCENT_DARK,
    )

    add_entity_box(
        slide,
        Inches(0.9),
        Inches(4.75),
        Inches(3.3),
        Inches(2.2),
        "Response",
        [
            "PK _id",
            "FK complaint -> Complaint._id",
            "FK respondedBy -> User._id",
            "message",
            "isInternal",
            "createdAt, updatedAt",
        ],
        title_color=GREEN,
    )

    add_entity_box(
        slide,
        Inches(8.9),
        Inches(4.75),
        Inches(3.7),
        Inches(2.2),
        "Notification",
        [
            "PK _id",
            "FK user -> User._id",
            "message, type",
            "FK relatedComplaint -> Complaint._id (nullable)",
            "isRead",
            "createdAt, updatedAt",
        ],
        title_color=AMBER,
    )

    # Relationship lines
    rel1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.9), Inches(2.55), Inches(4.8), Inches(2.55))
    rel1.line.color.rgb = MUTED
    rel1.line.width = Pt(1.5)
    t1 = slide.shapes.add_textbox(Inches(3.95), Inches(2.22), Inches(0.95), Inches(0.28))
    set_text_frame(t1.text_frame, "1..*", font_size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    t1b = slide.shapes.add_textbox(Inches(3.96), Inches(2.78), Inches(0.95), Inches(0.38))
    set_text_frame(t1b.text_frame, "submittedBy", font_size=9, bold=False, color=MUTED, align=PP_ALIGN.CENTER)

    rel2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(6.6), Inches(4.65), Inches(2.55), Inches(4.75))
    rel2.line.color.rgb = MUTED
    rel2.line.width = Pt(1.5)
    t2 = slide.shapes.add_textbox(Inches(4.2), Inches(4.55), Inches(1.2), Inches(0.3))
    set_text_frame(t2.text_frame, "1..*", font_size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

    rel3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.4), Inches(4.35), Inches(2.55), Inches(4.75))
    rel3.line.color.rgb = MUTED
    rel3.line.width = Pt(1.5)
    t3 = slide.shapes.add_textbox(Inches(1.62), Inches(4.44), Inches(0.9), Inches(0.3))
    set_text_frame(t3.text_frame, "1..*", font_size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    t3b = slide.shapes.add_textbox(Inches(0.98), Inches(4.42), Inches(0.8), Inches(0.3))
    set_text_frame(t3b.text_frame, "writes", font_size=9, bold=False, color=MUTED, align=PP_ALIGN.LEFT)

    rel4 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.5), Inches(2.95), Inches(9.0), Inches(5.15))
    rel4.line.color.rgb = MUTED
    rel4.line.width = Pt(1.5)
    t4 = slide.shapes.add_textbox(Inches(8.55), Inches(3.85), Inches(0.9), Inches(0.3))
    set_text_frame(t4.text_frame, "1..*", font_size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    t4b = slide.shapes.add_textbox(Inches(8.6), Inches(4.15), Inches(1.25), Inches(0.3))
    set_text_frame(t4b.text_frame, "relatedComplaint", font_size=8.5, bold=False, color=MUTED, align=PP_ALIGN.LEFT)

    rel5 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.9), Inches(3.5), Inches(8.9), Inches(5.15))
    rel5.line.color.rgb = RGBColor(148, 163, 184)
    rel5.line.width = Pt(1.2)
    t5 = slide.shapes.add_textbox(Inches(2.95), Inches(3.88), Inches(1.5), Inches(0.3))
    set_text_frame(t5.text_frame, "1..*", font_size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    t5b = slide.shapes.add_textbox(Inches(2.82), Inches(4.13), Inches(1.4), Inches(0.3))
    set_text_frame(t5b.text_frame, "receives", font_size=9, bold=False, color=MUTED, align=PP_ALIGN.CENTER)


def create_end_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(15, 23, 42)
    bg.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.9), Inches(1.4), Inches(11.8), Inches(1.0))
    set_text_frame(title.text_frame, "Project Summary", font_size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    body = slide.shapes.add_textbox(Inches(0.95), Inches(2.7), Inches(11.2), Inches(2.8))
    tf = body.text_frame
    tf.clear()
    lines = [
        "Layered MERN architecture with role-secure operations",
        "Complete complaint lifecycle: submit -> assign -> resolve -> notify",
        "Clear data model centered on User, Complaint, Response, and Notification",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.bullet = True
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(226, 232, 240)
        p.space_after = Pt(14)

    footer = slide.shapes.add_textbox(Inches(0.95), Inches(6.45), Inches(11.5), Inches(0.4))
    set_text_frame(
        footer.text_frame,
        "End of presentation",
        font_size=14,
        bold=False,
        color=RGBColor(148, 163, 184),
        align=PP_ALIGN.LEFT,
    )


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9
    prs.slide_height = Inches(7.5)

    create_title_slide(prs)
    create_architecture_slide(prs)
    create_module_slides(prs)
    create_data_flow_slide(prs)
    create_use_case_slide(prs)
    create_er_slide(prs)
    create_end_slide(prs)

    prs.save(OUTPUT_FILE)
    print(f"Created: {OUTPUT_FILE.resolve()}")


if __name__ == "__main__":
    build_presentation()
